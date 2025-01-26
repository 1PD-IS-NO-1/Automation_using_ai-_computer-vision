import os
import shutil
import streamlit as st
import cv2
import numpy as np
import fitz  # PyMuPDF
import pptx  # python-pptx
from cvzone.HandTrackingModule import HandDetector
import time 
def safe_remove_folder(output_folder):
    """Safely remove folder, handling permission issues"""
    try:
        if os.path.exists(output_folder):
            # Remove read-only files
            for root, dirs, files in os.walk(output_folder):
                for file in files:
                    file_path = os.path.join(root, file)
                    os.chmod(file_path, 0o666)
            shutil.rmtree(output_folder)
    except Exception as e:
        st.warning(f"Could not completely clear folder: {e}")
    
    # Recreate the folder
    os.makedirs(output_folder, exist_ok=True)

def extract_images_from_pdf(pdf_path, output_folder):
    # Safely remove existing folder
  

    # Open PDF
    pdf_document = fitz.open(pdf_path)
    
    # Convert pages to images
    for page_num in range(len(pdf_document)):
        page = pdf_document[page_num]
        pix = page.get_pixmap(dpi=300)  # Higher resolution
        img_path = os.path.join(output_folder, f"page_{page_num + 1}.jpg")
        pix.save(img_path)

    return len(pdf_document)

def extract_images_from_pptx(pptx_path, output_folder):
    # Safely remove existing folder
    

    # Open PowerPoint
    presentation = pptx.Presentation(pptx_path)
    
    # Convert slides to images
    for slide_num, slide in enumerate(presentation.slides):
        slide_image_path = os.path.join(output_folder, f"slide_{slide_num + 1}.jpg")
        
        # Create blank white image with actual slide dimensions
        slide_width = int(presentation.slide_width / 12700)  # Convert to pixels
        slide_height = int(presentation.slide_height / 12700)  # Convert to pixels
        slide_image = np.ones((slide_height, slide_width, 3), dtype=np.uint8) * 255
        
        # Save slide image
        cv2.imwrite(slide_image_path, slide_image)

    return len(presentation.slides)

def hand_slide_navigation(folderPath):
    # Parameters
    width, height = 880, 420
    gestureThreshold = 300

    # Camera Setup
    cap = cv2.VideoCapture(0)
    cap.set(3, width)
    cap.set(4, height)

    # Hand Detector
    detectorHand = HandDetector(detectionCon=0.8, maxHands=1)

    # Variables
    imgNumber = 0
    delay = 30
    counter = 0
    buttonPressed = False
    annotations = [[]]
    annotationNumber = -1
    annotationStart = False
    hs, ws = int(120 * 1), int(213 * 1)

    # Get list of presentation images
    pathImages = sorted(os.listdir(folderPath), key=len)
    
    # Tracking variables for hand movement
    prev_hand_center = None
    movement_threshold = 50

    while True:
        # Get image frame
        success, img = cap.read()
        img = cv2.flip(img, 1)
        pathFullImage = os.path.join(folderPath, pathImages[imgNumber])
        imgCurrent = cv2.imread(pathFullImage)

        # Find the hand and its landmarks
        hands, img = detectorHand.findHands(img)

        # Draw Gesture Threshold line
        cv2.line(img, (0, gestureThreshold), (width, gestureThreshold), (0, 255, 0), 10)

        if hands and not buttonPressed:
            hand = hands[0]
            cx, cy = hand["center"]
            lmList = hand["lmList"]
            fingers = detectorHand.fingersUp(hand)

            # Check for full hand movement (all fingers extended)
            if all(finger == 1 for finger in fingers):
                # First time detecting hand or enough movement detected
                if prev_hand_center is None:
                    prev_hand_center = cx
                else:
                    # Detect horizontal movement
                    movement = cx - prev_hand_center
                    
                    if abs(movement) > movement_threshold:
                        if movement > 0:  # Moving right
                            if imgNumber < len(pathImages) - 1:
                                imgNumber += 1
                                annotations = [[]]
                                annotationNumber = -1
                                buttonPressed = True
                        else:  # Moving left
                            if imgNumber > 0:
                                imgNumber -= 1
                                annotations = [[]]
                                annotationNumber = -1
                                buttonPressed = True
                    
                    # Update previous center
                    prev_hand_center = cx

            # Annotation functionality
            if fingers == [0, 1, 0, 0, 0]:
                xVal = int(np.interp(lmList[8][0], [width // 2, width], [0, width]))
                yVal = int(np.interp(lmList[8][1], [150, height-150], [0, height]))
                indexFinger = xVal, yVal

                if annotationStart is False:
                    annotationStart = True
                    annotationNumber += 1
                    annotations.append([])

                annotations[annotationNumber].append(indexFinger)
                cv2.circle(imgCurrent, indexFinger, 12, (0, 0, 255), cv2.FILLED)
            else:
                annotationStart = False

        # Slide change delay mechanism
        if buttonPressed:
            counter += 1
            if counter > delay:
                counter = 0
                buttonPressed = False
                prev_hand_center = None

        # Draw annotations
        for annotation in annotations:
            for j in range(1, len(annotation)):
                cv2.line(imgCurrent, annotation[j-1], annotation[j], (0, 0, 200), 12)

        # Add small camera view to slide
        imgSmall = cv2.resize(img, (ws, hs))
        h, w, _ = imgCurrent.shape
        imgCurrent[0:hs, w - ws: w] = imgSmall

        cv2.imshow("Slides", imgCurrent)
        cv2.imshow("Image", img)

        key = cv2.waitKey(1)
        if key == ord('q'):
            break

    cap.release()
    cv2.destroyAllWindows()

def main():
    st.title("Slide Navigator with Hand Tracking")

    # File Upload
    uploaded_file = st.file_uploader("Choose a PDF or PPTX file", type=["pdf", "pptx"])
    
    if uploaded_file is not None:
        # Create a temporary directory for slides
        temp_folder = "Presentation"
        
        # Ensure the temporary folder exists
        os.makedirs(temp_folder, exist_ok=True)
        safe_remove_folder(temp_folder)

        # Save uploaded file
        file_path = os.path.join(temp_folder, uploaded_file.name)
        with open(file_path, "wb") as f:
            f.write(uploaded_file.getvalue())
        
        # Extract images based on file type
        try:
            if uploaded_file.name.lower().endswith('.pdf'):
                total_pages = extract_images_from_pdf(file_path, temp_folder)
            elif uploaded_file.name.lower().endswith('.pptx'):
                total_pages = extract_images_from_pptx(file_path, temp_folder)
            
            st.success(f"Slides Extracted: {total_pages} pages")
            
            # Start Hand Navigation Button
            if st.button("Start Hand Navigation"):
                hand_slide_navigation(temp_folder)
        
        except Exception as e:
            st.error(f"Error processing file: {e}")

if __name__ == "__main__":
    main()